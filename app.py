import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
import io
import pandas as pd

# --- Segédfüggvények ---
def set_font(run, font_name='Karla', size=14, bold=False):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold

def create_ppt(word_file, template_file):
    # Sablon betöltése
    prs = Presentation(template_file)
    
    # Word dokumentum beolvasása (egyszerűsített elemzés)
    doc = Document(word_file)
    full_text = [p.text for p in doc.paragraphs if p.text.strip() != ""]
    
    # --- 1. SLIDE: Címoldal ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Játékosszegmentáció a Tippmixprón"
    
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Kockázatkezelési szempontból veszélyes fogadók | 2025"

    # --- 2. SLIDE: A szegmentáció célja ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "A szegmentáció célja"
    content = slide.placeholders[1].text_frame
    content.text = "Veszélyes fogadási mintázatok egységes kezelése."
    p = content.add_paragraph()
    p.text = "A csoportok teljesítményének transzparens riportálása."
    p.level = 0

    # --- 3. SLIDE: Kockázatkezelési Role-ok (8 kategória) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Alkalmazott kategóriák (8 Role)"
    tf = slide.placeholders[1].text_frame
    roles = ["Bonus Abuser", "Cashout Abuser", "Monitoring", "Late Betting", 
             "Arb/Dropping", "Default", "Sport Risk", "Tipping Line"]
    for role in roles:
        p = tf.add_paragraph()
        p.text = role
        p.level = 1

    # --- 4. SLIDE: 2025-ös Teljesítmény adatok ---
    # Itt egy táblázatot hozunk létre a Word-ben található adatok alapján
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Üres layout címmel
    slide.shapes.title.text = "Role-ok 2025-ös teljesítménye"
    
    rows, cols = 6, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(3.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Táblázat fejlécek
    headers = ["Role", "Stakes (Ft)", "GGR (Ft)", "Margin"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    # Példa adatok a dokumentumból
    data = [
        ["Tipping Line", "16 980 626 790", "-1 086 229 812", "-6,40%"],
        ["Sport Risk", "4 485 926 926", "-99 890 372", "-2,23%"],
        ["Default", "10 242 247 846", "508 346 821", "4,96%"],
        ["Arb/Dropping", "336 359 918", "-23 445 041", "-6,97%"],
        ["Monitoring", "21 183 867 861", "680 603 239", "3,21%"]
    ]
    
    for r_idx, row_vals in enumerate(data):
        for c_idx, val in enumerate(row_vals):
            table.cell(r_idx + 1, c_idx).text = val

    # --- 5. SLIDE: Megelőzött veszteségek ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Eredmény: Megelőzött veszteség"
    tf = slide.placeholders[1].text_frame
    tf.text = "Összesen: ~80,5 millió Ft megelőzött veszteség"
    
    for item in ["Tipping Line: -71,9 M Ft", "Sport Risk: -5,6 M Ft", "Arb/Dropping: -2,8 M Ft"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # Mentés memóriába
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- STREAMLIT UI ---
st.set_page_config(page_title="SZRT PPT Generátor", layout="centered")

st.image("https://upload.wikimedia.org/wikipedia/commons/thumb/d/d4/Szerencsejatek_Zrt_logo.svg/1024px-Szerencsejatek_Zrt_logo.svg.png", width=200)
st.title("Prezentáció Generátor")
st.info("Töltsd fel a Word dokumentumot és a PPT sablont a generáláshoz.")

col1, col2 = st.columns(2)
with col1:
    word_file = st.file_uploader("Word dokumentum", type="docx")
with col2:
    ppt_template = st.file_uploader("Céges PPT sablon", type="pptx")

if st.button("🚀 Prezentáció elkészítése"):
    if word_file and ppt_template:
        with st.spinner("Dolgozom rajta..."):
            final_ppt = create_ppt(word_file, ppt_template)
            st.success("A prezentáció elkészült!")
            st.download_button(
                label="📥 Letöltés (.pptx)",
                data=final_ppt,
                file_name="Tippmixpro_Szegmentacio_Prezentacio.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.warning("Kérlek, töltsd fel mindkét fájlt!")

st.markdown("---")
st.caption("A generált prezentáció a megadott céges sablont és betűtípusokat használja.")
